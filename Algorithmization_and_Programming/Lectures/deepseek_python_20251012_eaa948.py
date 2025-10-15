# pptx_to_md_advanced.py
from pptx import Presentation
import re

class PPTXConverter:
    def __init__(self):
        self.code_blocks = []
        self.current_slide = 0
    
    def convert(self, pptx_path, md_path):
        prs = Presentation(pptx_path)
        
        with open(md_path, 'w', encoding='utf-8') as md_file:
            md_file.write("# Презентация\n\n")
            
            for slide_num, slide in enumerate(prs.slides, 1):
                self.current_slide = slide_num
                md_file.write(f"## Слайд {slide_num}\n\n")
                
                self.process_slide(slide, md_file)
                md_file.write("\n---\n\n")
        
        print(f"✅ Успешно конвертировано {len(prs.slides)} слайдов")
    
    def process_slide(self, slide, md_file):
        # Обработка заголовка
        if slide.shapes.title:
            title = self.clean_text(slide.shapes.title.text)
            if title:
                md_file.write(f"### {title}\n\n")
        
        # Обработка содержимого
        text_elements = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                if shape == slide.shapes.title:
                    continue
                
                text = self.clean_text(shape.text)
                if text:
                    text_elements.append(text)
        
        # Форматирование вывода
        for text in text_elements:
            if self.is_code_block(text):
                md_file.write("```cpp\n")
                md_file.write(text + "\n")
                md_file.write("```\n\n")
            else:
                # Обработка обычного текста
                lines = text.split('\n')
                for line in lines:
                    line = line.strip()
                    if line:
                        md_file.write(line + "\n\n")
    
    def clean_text(self, text):
        """Очистка текста от лишних пробелов"""
        text = re.sub(r'\s+', ' ', text)
        return text.strip()
    
    def is_code_block(self, text):
        """Определение кода C++"""
        code_patterns = [
            r'#include\s*<.*>',
            r'int main\s*\(\)',
            r'void\s+\w+\s*\(',
            r'cout\s*<<',
            r'printf\s*\(',
            r'using namespace',
            r'struct\s+\w+',
            r'class\s+\w+'
        ]
        
        text_first_line = text.split('\n')[0].strip()
        return any(re.search(pattern, text_first_line, re.IGNORECASE) 
                  for pattern in code_patterns)

# Использование
if __name__ == "__main__":
    converter = PPTXConverter()
    converter.convert("АИП_1_3.pptx", "презентация_advanced.md")